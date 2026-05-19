#!/usr/bin/env python3
"""
Extract every shape on every slide of a .pptx into a JSON dump.

Usage:
    python extract_pptx.py <path/to/deck.pptx>            # prints to stdout
    python extract_pptx.py <path/to/deck.pptx> -o dump.json
    python extract_pptx.py <path/to/deck.pptx> -o dump.json --assets-dir assets/

The dump captures the *actual* state of the export — text content, position,
size, and per-run typography (font name, size, bold, italic, color). Use this
as the ground truth for the fidelity audit; do not trust the export script's
intent.

For the reconstruction workflow (rebuilding an externally-authored deck that
has no HTML source), the dump also carries the content a faithful redesign
must preserve: per-slide speaker `notes`, per-shape `table` cell text, and
per-shape `image` metadata. Pass `--assets-dir <dir>` to also write the image
blobs to disk so the rebuilt deck can re-embed logos/charts. All of these are
*additive* — existing keys are unchanged, so the audit workflow and the
`extract_pptx.py deck.pptx > pptx_dump.json` invocation behave exactly as
before.

Coordinates are reported in inches (rounded to 3 decimals) so they're
human-readable when comparing against rails like CONTENT_MAX_Y = 6.70".
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.stderr.write(
        "python-pptx is required. Install with: pip install python-pptx\n"
    )
    sys.exit(2)


def emu_to_in(emu: int | None) -> float | None:
    if emu is None:
        return None
    return round(emu / 914400, 3)


def color_repr(color) -> str | None:
    """Best-effort color extraction. Returns hex string or None."""
    if color is None:
        return None
    try:
        # ColorFormat.type may be None when no explicit color is set.
        if color.type is None:
            return None
        rgb = color.rgb
        if rgb is None:
            return None
        return f"#{str(rgb).lower()}"
    except (AttributeError, ValueError, TypeError):
        return None


def extract_runs(text_frame) -> list[dict]:
    runs = []
    for para in text_frame.paragraphs:
        for run in para.runs:
            font = run.font
            runs.append({
                "text": run.text,
                "font": font.name,
                "size_pt": float(font.size.pt) if font.size is not None else None,
                "bold": bool(font.bold) if font.bold is not None else None,
                "italic": bool(font.italic) if font.italic is not None else None,
                # Color is independent of font name/size: a run can inherit
                # font from the theme yet set its own color. Color drift is
                # one of the things this audit needs to catch, so don't gate
                # the extraction on unrelated font attributes.
                "color": color_repr(font.color),
            })
    return runs


def extract_table(shape) -> dict:
    """Cell text as a row-major matrix. Geometry stays on the shape dict."""
    table = shape.table
    rows = []
    for row in table.rows:
        rows.append([cell.text for cell in row.cells])
    return {"rows": rows}


def extract_image(shape, *, assets_dir: Path | None,
                  slide_index: int, shape_index: int) -> dict:
    """Image metadata; optionally write the blob to assets_dir.

    When assets_dir is None the blob is not written and `filename` is null,
    but the dump still records that an image occupies this shape so the
    reconstruction can decide whether it needs the asset re-extracted.
    """
    img = shape.image
    info: dict = {"content_type": img.content_type, "filename": None}
    if assets_dir is not None:
        ext = (img.ext or "bin").lstrip(".")
        name = f"slide-{slide_index:03d}-img-{shape_index:03d}.{ext}"
        assets_dir.mkdir(parents=True, exist_ok=True)
        (assets_dir / name).write_bytes(img.blob)
        info["filename"] = name
    return info


def extract_shape(shape, *, assets_dir: Path | None = None,
                   slide_index: int = 0, shape_index: int = 0) -> dict:
    data = {
        "name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type is not None else None,
        "left_in": emu_to_in(shape.left),
        "top_in": emu_to_in(shape.top),
        "width_in": emu_to_in(shape.width),
        "height_in": emu_to_in(shape.height),
    }
    if shape.left is not None and shape.height is not None and shape.top is not None:
        data["bottom_in"] = emu_to_in(shape.top + shape.height)
        data["right_in"] = emu_to_in(shape.left + shape.width)
    if shape.has_text_frame:
        tf = shape.text_frame
        data["text"] = tf.text
        data["runs"] = extract_runs(tf)
    if shape.has_table:
        data["table"] = extract_table(shape)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            data["image"] = extract_image(
                shape, assets_dir=assets_dir,
                slide_index=slide_index, shape_index=shape_index,
            )
        except (AttributeError, KeyError, ValueError):
            # Linked/placeholder pictures with no embedded blob: record the
            # slot so the redesign knows an image was here, even if we can't
            # recover the bytes.
            data["image"] = {"content_type": None, "filename": None}
    return data


def slide_notes(slide) -> str | None:
    if not slide.has_notes_slide:
        return None
    text = slide.notes_slide.notes_text_frame.text
    return text if text.strip() else None


def extract_pptx(path: Path, assets_dir: Path | None = None) -> dict:
    prs = Presentation(str(path))
    canvas = {
        "width_in": emu_to_in(prs.slide_width),
        "height_in": emu_to_in(prs.slide_height),
    }
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        shapes = [
            extract_shape(s, assets_dir=assets_dir, slide_index=i, shape_index=j)
            for j, s in enumerate(slide.shapes, 1)
        ]
        slides.append({
            "index": i,
            "notes": slide_notes(slide),
            "shapes": shapes,
        })
    return {
        "source": str(path),
        "canvas": canvas,
        "slide_count": len(slides),
        "slides": slides,
    }


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    ap.add_argument("path", type=Path, help=".pptx file to extract")
    ap.add_argument("-o", "--output", type=Path, help="write JSON to this path; default stdout")
    ap.add_argument("--assets-dir", type=Path, default=None,
                    help="if set, write embedded image blobs into this directory "
                         "and record their filenames in the dump (for the "
                         "reconstruction workflow). Default: do not extract blobs.")
    args = ap.parse_args()

    if not args.path.exists():
        ap.error(f"file not found: {args.path}")

    data = extract_pptx(args.path, assets_dir=args.assets_dir)
    payload = json.dumps(data, ensure_ascii=False, indent=2)
    if args.output:
        args.output.write_text(payload, encoding="utf-8")
        sys.stderr.write(f"wrote {args.output} ({len(payload)} bytes, {data['slide_count']} slides)\n")
    else:
        sys.stdout.write(payload)
        sys.stdout.write("\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
